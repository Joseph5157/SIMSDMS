import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_path):
    prs = Presentation()
    # Set to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    c_dark_navy = RGBColor(15, 23, 42)      # #0f172a
    c_brand_blue = RGBColor(37, 99, 235)     # #2563eb
    c_light_slate = RGBColor(248, 250, 252)  # #f8fafc
    c_border_gray = RGBColor(226, 232, 240)  # #e2e8f0
    c_text_ink = RGBColor(30, 41, 59)        # #1e293b
    c_text_muted = RGBColor(100, 116, 139)   # #64748b
    c_white = RGBColor(255, 255, 255)
    c_accent_green = RGBColor(16, 185, 129)  # #10b981
    c_accent_amber = RGBColor(245, 158, 11)  # #f59e0b
    c_accent_red = RGBColor(239, 68, 68)     # #ef4444

    blank_layout = prs.slide_layouts[6] # completely blank layout

    # Helper function to add a header to standard slides
    def add_slide_header(slide, title_text, eyebrow_text):
        # Eyebrow
        eb_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        eb_tf = eb_box.text_frame
        eb_tf.word_wrap = True
        eb_tf.margin_left = eb_tf.margin_top = eb_tf.margin_right = eb_tf.margin_bottom = 0
        eb_p = eb_tf.paragraphs[0]
        eb_p.text = eyebrow_text.upper()
        eb_p.font.name = "Arial"
        eb_p.font.size = Pt(11)
        eb_p.font.bold = True
        eb_p.font.color.rgb = c_brand_blue

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_top = title_tf.margin_right = title_tf.margin_bottom = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = "Arial"
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = c_dark_navy

    # Helper function to set solid background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to add standard speaker notes
    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    # ---------------------------------------------------------
    # SLIDE 1: Cover (Dark Navy Theme)
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_dark_navy)

    # Accent decorative block
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c_brand_blue
    shape.line.fill.background()

    # Cover Title Box
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_eb = tf.paragraphs[0]
    p_eb.text = "SIMS COLLEGE OF PHARMACY · DISCIPLINE MANAGEMENT"
    p_eb.font.name = "Arial"
    p_eb.font.size = Pt(12)
    p_eb.font.bold = True
    p_eb.font.color.rgb = c_brand_blue
    p_eb.space_after = Pt(14)

    p_title = tf.add_paragraph()
    p_title.text = "Your Discipline Duty,\nOn Your Phone"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(54)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    p_title.space_after = Pt(20)

    p_sub = tf.add_paragraph()
    p_sub.text = "A quick, hands-on faculty guide to using the SIMS DMS PWA application.\nNo paper files. No chasing updates. All in one place."
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = c_text_muted

    # Key highlight pills
    pills_data = ["📲 Works on Mobile", "✈️ OTP via Telegram", "⏱️ 5 Core Tasks"]
    pill_left = Inches(1.2)
    for pill in pills_data:
        pill_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_left, Inches(5.2), Inches(2.2), Inches(0.45))
        pill_shape.fill.solid()
        pill_shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
        pill_shape.line.color.rgb = c_brand_blue
        pill_shape.line.width = Pt(1)
        
        ptf = pill_shape.text_frame
        ptf.word_wrap = True
        ptf.margin_top = ptf.margin_bottom = ptf.margin_left = ptf.margin_right = Inches(0.05)
        p = ptf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = pill
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = c_white
        pill_left += Inches(2.4)

    add_notes(slide, 
        "Welcome slide. Introduces the transition to the mobile-first SIMS DMS (Discipline Management System) at SIMS College of Pharmacy.\n\n"
        "Key points to cover:\n"
        "1. Welcome the faculty to the training session.\n"
        "2. Explain that the main goal of this system is to replace manual paper logs at the gate and WhatsApp group coordinates.\n"
        "3. Emphasize that it works directly from any browser on their mobile phone, and uses Telegram for login OTPs and notifications.")

    # ---------------------------------------------------------
    # SLIDE 2: Why DMS (Comparison / Value Prop)
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_light_slate)
    add_slide_header(slide, "Why We're Moving to SIMS DMS", "The Benefits")

    # Column 1: Problems before
    col1_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    col1_bg.fill.solid()
    col1_bg.fill.fore_color.rgb = c_white
    col1_bg.line.color.rgb = c_border_gray
    col1_bg.line.width = Pt(1)
    
    tf1 = col1_bg.text_frame
    tf1.word_wrap = True
    tf1.margin_top = tf1.margin_left = tf1.margin_right = Inches(0.4)
    
    p = tf1.paragraphs[0]
    p.text = "BEFORE (Manual System)"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = c_accent_red
    p.space_after = Pt(14)

    bullets_before = [
        "❌ Paper register at the gate (hard to sign, verify, and track).",
        "❌ Duty coverage / swaps arranged over chaotic phone calls and WhatsApp.",
        "❌ Student violations noted in books (no shared status with the office).",
        "❌ 'Were you on duty?' disputes due to lack of tamper-proof logs."
    ]
    for b in bullets_before:
        bp = tf1.add_paragraph()
        bp.text = b
        bp.font.name = "Arial"
        bp.font.size = Pt(13)
        bp.font.color.rgb = c_text_ink
        bp.space_after = Pt(12)

    # Column 2: Solutions after
    col2_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8))
    col2_bg.fill.solid()
    col2_bg.fill.fore_color.rgb = RGBColor(239, 246, 255) # light blue
    col2_bg.line.color.rgb = c_brand_blue
    col2_bg.line.width = Pt(1.5)
    
    tf2 = col2_bg.text_frame
    tf2.word_wrap = True
    tf2.margin_top = tf2.margin_left = tf2.margin_right = Inches(0.4)
    
    p = tf2.paragraphs[0]
    p.text = "WITH SIMS DMS (Digital System)"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = c_brand_blue
    p.space_after = Pt(14)

    bullets_after = [
        "✅ Clean, automated check-in and check-out timestamped in the app.",
        "✅ Official swaps sent and accepted in 2 taps with full digital tracking.",
        "✅ Fines auto-calculated by system settings and saved instantly.",
        "✅ Proactive Telegram bot notifications for schedules and reminders."
    ]
    for b in bullets_after:
        bp = tf2.add_paragraph()
        bp.text = b
        bp.font.name = "Arial"
        bp.font.size = Pt(13)
        bp.font.color.rgb = c_text_ink
        bp.space_after = Pt(12)

    add_notes(slide,
        "Speaker Notes: Contrast the old process with the new workflow.\n\n"
        "Key points:\n"
        "1. Mention the friction in the past manual system (lost registers, undocumented duty switches, manual billing mistakes).\n"
        "2. Detail how DMS automates checking in/out, reassignments (swaps), and violation recording.\n"
        "3. Emphasize that the Telegram Bot acts as an active assistant rather than forcing faculty to constantly log in to check for updates.")

    # ---------------------------------------------------------
    # SLIDE 3: Getting Started & Navigation
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_light_slate)
    add_slide_header(slide, "Signing In & Navigating the System", "Orientation")

    # Step list
    steps_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    stf = steps_box.text_frame
    stf.word_wrap = True
    stf.margin_left = stf.margin_top = stf.margin_right = stf.margin_bottom = 0
    
    p = stf.paragraphs[0]
    p.text = "How to log in:"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark_navy
    p.space_after = Pt(10)

    steps = [
        "1. Open the app link on your browser: input your registered email address.",
        "2. Click 'Request OTP'. The system will generate a secure, 6-digit code.",
        "3. Check your Telegram app. The official SIMS DMS Bot will deliver the code.",
        "4. Enter the code in the app to authenticate. No passwords to memorize!",
        "5. If you haven't linked Telegram yet, use the activation link sent by your admin."
    ]
    for step in steps:
        sp = stf.add_paragraph()
        sp.text = step
        sp.font.name = "Arial"
        sp.font.size = Pt(13.5)
        sp.font.color.rgb = c_text_ink
        sp.space_after = Pt(12)
        
    callout_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.8), Inches(6.5), Inches(1.2))
    callout_box.fill.solid()
    callout_box.fill.fore_color.rgb = RGBColor(254, 243, 199) # light amber
    callout_box.line.color.rgb = c_accent_amber
    callout_box.line.width = Pt(1)
    ctf = callout_box.text_frame
    ctf.word_wrap = True
    ctf.margin_top = ctf.margin_bottom = ctf.margin_left = ctf.margin_right = Inches(0.15)
    cp = ctf.paragraphs[0]
    cp.text = "💡 NOTE: The login session is secure and cached. You will only need to re-authenticate with an OTP after logouts or when cookies expire."
    cp.font.name = "Arial"
    cp.font.size = Pt(11)
    cp.font.color.rgb = RGBColor(146, 64, 14) # dark amber

    # Navigation illustration (Right Column)
    nav_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.8))
    nav_bg.fill.solid()
    nav_bg.fill.fore_color.rgb = c_white
    nav_bg.line.color.rgb = c_border_gray
    nav_bg.line.width = Pt(1)
    
    ntf = nav_bg.text_frame
    ntf.word_wrap = True
    ntf.margin_top = ntf.margin_left = ntf.margin_right = Inches(0.3)
    
    np = ntf.paragraphs[0]
    np.text = "THE FIVE BOTTOM TABS:"
    np.font.name = "Arial"
    np.font.size = Pt(14)
    np.font.bold = True
    np.font.color.rgb = c_dark_navy
    np.space_after = Pt(14)

    tabs = [
        "🏠 Dashboard: Check-ins, today's schedule, alerts.",
        "📅 Slots: Interactive calendar to pick monthly slots.",
        "✅ Attendance: Personal check-in logs and status summary.",
        "⚠️ Violations: List of student infractions you recorded.",
        "✉️ Messages: Official administrative announcements."
    ]
    for tab in tabs:
        tp = ntf.add_paragraph()
        tp.text = tab
        tp.font.name = "Arial"
        tp.font.size = Pt(13)
        tp.font.color.rgb = c_text_ink
        tp.space_after = Pt(12)

    add_notes(slide,
        "Speaker Notes: Describe how to get into the app and get oriented.\n\n"
        "Key points:\n"
        "1. Confirm that password resets and credential locks are a thing of the past. Login requires only email and a code delivered via Telegram.\n"
        "2. Explain the prerequisite: The Telegram account must be linked. Admins send a direct setup link or users can start the bot (`/start invite_token`).\n"
        "3. Walk through the bottom tabs. Reiterate that the app acts as a Mobile PWA, meaning they can add it to their home screen for easy access.")

    # ---------------------------------------------------------
    # SLIDE 4: Task 1 - Check In / Out
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_light_slate)
    add_slide_header(slide, "Task 1: Checking In & Out of Duty", "Daily Routine")

    # Text instructions
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    ttf = text_box.text_frame
    ttf.word_wrap = True
    ttf.margin_left = ttf.margin_top = ttf.margin_right = ttf.margin_bottom = 0
    
    p = ttf.paragraphs[0]
    p.text = "How it works on duty days:"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark_navy
    p.space_after = Pt(12)

    bullets_t1 = [
        "1. Open the app on the Dashboard. Today's assigned duty will show as a blue 'Today's Duty' card at the top.",
        "2. When you arrive at your designated gate or duty area, tap the 'Check In' button.",
        "3. When your shift ends and you are leaving, open the app again and tap 'Check Out'.",
        "4. If you have multiple sessions on the same day (e.g., Morning and Afternoon), each session has its own card. Check in/out for each session separately."
    ]
    for b in bullets_t1:
        bp = ttf.add_paragraph()
        bp.text = b
        bp.font.name = "Arial"
        bp.font.size = Pt(13.5)
        bp.font.color.rgb = c_text_ink
        bp.space_after = Pt(12)

    # Auto Clock-Out Warning Callout
    warn_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.0), Inches(6.5), Inches(1.2))
    warn_box.fill.solid()
    warn_box.fill.fore_color.rgb = RGBColor(254, 242, 242) # light red
    warn_box.line.color.rgb = c_accent_red
    warn_box.line.width = Pt(1)
    wtf = warn_box.text_frame
    wtf.word_wrap = True
    wtf.margin_top = wtf.margin_bottom = wtf.margin_left = wtf.margin_right = Inches(0.15)
    wp = wtf.paragraphs[0]
    wp.text = "⚠️ Auto-checkout warning:\nIf you forget to check out, the system will automatically check you out at the cutoff hour. This logs as 'Auto' in your attendance and triggers a warning report for the admin. Please check out manually!"
    wp.font.name = "Arial"
    wp.font.size = Pt(11)
    wp.font.color.rgb = RGBColor(153, 27, 27) # dark red

    # Graphic representation (Mock Card)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(3.2))
    card.fill.solid()
    card.fill.fore_color.rgb = c_brand_blue
    card.line.fill.background()
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.margin_top = ctf.margin_left = ctf.margin_right = Inches(0.3)
    
    cp1 = ctf.paragraphs[0]
    cp1.text = "📋 TODAY'S DUTY"
    cp1.font.name = "Arial"
    cp1.font.size = Pt(11)
    cp1.font.bold = True
    cp1.font.color.rgb = RGBColor(191, 219, 254) # light blue-200
    cp1.space_after = Pt(6)

    cp2 = ctf.add_paragraph()
    cp2.text = "Morning Session"
    cp2.font.name = "Arial"
    cp2.font.size = Pt(24)
    cp2.font.bold = True
    cp2.font.color.rgb = c_white
    cp2.space_after = Pt(4)

    cp3 = ctf.add_paragraph()
    cp3.text = "Starts 8:00 AM · Not Checked In"
    cp3.font.name = "Arial"
    cp3.font.size = Pt(13)
    cp3.font.color.rgb = RGBColor(219, 234, 254)
    cp3.space_after = Pt(18)

    # Button inside card graphic
    btn_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(3.8), Inches(3.9), Inches(0.55))
    btn_shape.fill.solid()
    btn_shape.fill.fore_color.rgb = c_white
    btn_shape.line.fill.background()
    btf = btn_shape.text_frame
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.text = "Tap to Check In"
    bp.font.name = "Arial"
    bp.font.size = Pt(14)
    bp.font.bold = True
    bp.font.color.rgb = c_brand_blue

    add_notes(slide,
        "Speaker Notes: Walk through the core daily routine of checking in and out.\n\n"
        "Key points:\n"
        "1. Stress that checking in is mandatory. It verifies physical presence on campus for disciplinary audits.\n"
        "2. Explain that the app uses background geo-validation or simply relies on honest timestamp reporting.\n"
        "3. Highlight the auto clock-out feature. While it prevents cards from staying open indefinitely, having too many 'Auto' checks on record reflects poorly on the faculty attendance report.")

    # ---------------------------------------------------------
    # SLIDE 5: Task 2 - Picking Slots
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_light_slate)
    add_slide_header(slide, "Task 2: Picking Your Monthly Slots", "Start of the Month")

    # Left Column Text
    slots_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    s_tf = slots_text.text_frame
    s_tf.word_wrap = True
    s_tf.margin_left = s_tf.margin_top = s_tf.margin_right = s_tf.margin_bottom = 0
    
    p = s_tf.paragraphs[0]
    p.text = "How to select your dates:"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark_navy
    p.space_after = Pt(12)

    bullets_t2 = [
        "1. Go to the 'Slots' tab. A green banner will display indicating that the scheduling window is open.",
        "2. The calendar will highlight dates. A blue dot indicates that duty slots are still available for selection.",
        "3. Tap on your desired date. A selection drawer will slide up from the bottom.",
        "4. Choose 'Pick Morning' or 'Pick Afternoon'. Once confirmed, the date turns green.",
        "5. Repeat this process until you meet your required session target count for the month (typically 3 slots)."
    ]
    for b in bullets_t2:
        bp = s_tf.add_paragraph()
        bp.text = b
        bp.font.name = "Arial"
        bp.font.size = Pt(13.5)
        bp.font.color.rgb = c_text_ink
        bp.space_after = Pt(12)

    # Info banner
    info_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.2), Inches(6.5), Inches(1.0))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = RGBColor(239, 246, 255) # light blue
    info_box.line.color.rgb = c_brand_blue
    info_box.line.width = Pt(1)
    itf = info_box.text_frame
    itf.word_wrap = True
    itf.margin_top = itf.margin_bottom = itf.margin_left = itf.margin_right = Inches(0.15)
    ip = itf.paragraphs[0]
    ip.text = "💡 Did you make a mistake?\nOnce a slot is picked, you cannot drop it directly in the app. You must request a swap with a colleague (Task 4) or contact an Admin to clear it."
    ip.font.name = "Arial"
    ip.font.size = Pt(11)
    ip.font.color.rgb = c_brand_blue

    # Calendar Graphic representation
    cal_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.8))
    cal_box.fill.solid()
    cal_box.fill.fore_color.rgb = c_white
    cal_box.line.color.rgb = c_border_gray
    cal_box.line.width = Pt(1)
    ctf = cal_box.text_frame
    ctf.word_wrap = True
    ctf.margin_top = ctf.margin_left = ctf.margin_right = Inches(0.3)
    
    cp = ctf.paragraphs[0]
    cp.text = "SLOT PICKER CALENDAR"
    cp.font.name = "Arial"
    cp.font.size = Pt(12)
    cp.font.bold = True
    cp.font.color.rgb = c_text_muted
    cp.space_after = Pt(20)

    # Let's draw a mock calendar grid
    grid_top = Inches(2.5)
    cell_size = Inches(0.7)
    gap = Inches(0.1)
    days = [
        ["1", "2", "3", "4", "5", "6", "7"],
        ["8", "9", "10", "11", "12", "13", "14"]
    ]
    # Highlight codes: 0 = regular, 1 = available (blue outline), 2 = picked (green fill)
    highlights = [
        [0, 1, 0, 2, 1, 0, 0],
        [0, 1, 1, 0, 1, 0, 0]
    ]

    for r_idx, row in enumerate(days):
        for c_idx, day_num in enumerate(row):
            left = Inches(8.0) + c_idx * (cell_size + gap)
            top = grid_top + r_idx * (cell_size + gap)
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, cell_size, cell_size)
            
            h = highlights[r_idx][c_idx]
            if h == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
                cell.line.color.rgb = c_border_gray
            elif h == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(239, 246, 255)
                cell.line.color.rgb = c_brand_blue
                cell.line.width = Pt(1.5)
            elif h == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(209, 250, 229) # light green
                cell.line.color.rgb = c_accent_green
                cell.line.width = Pt(1.5)
                
            cell_tf = cell.text_frame
            cell_p = cell_tf.paragraphs[0]
            cell_p.alignment = PP_ALIGN.CENTER
            cell_p.text = day_num
            cell_p.font.name = "Arial"
            cell_p.font.size = Pt(12)
            cell_p.font.bold = True
            
            if h == 0:
                cell_p.font.color.rgb = c_text_muted
            elif h == 1:
                cell_p.font.color.rgb = c_brand_blue
            elif h == 2:
                cell_p.font.color.rgb = c_accent_green

    # Pick button simulation below grid
    btn_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(4.3), Inches(4.3), Inches(0.5))
    btn_shape.fill.solid()
    btn_shape.fill.fore_color.rgb = c_brand_blue
    btn_shape.line.fill.background()
    btf = btn_shape.text_frame
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.text = "Pick Morning Session"
    bp.font.name = "Arial"
    bp.font.size = Pt(12)
    bp.font.bold = True
    bp.font.color.rgb = c_white

    add_notes(slide,
        "Speaker Notes: Explain how scheduling works.\n\n"
        "Key points:\n"
        "1. The admin will announce when the monthly scheduling window opens (via Telegram broadcast).\n"
        "2. Demonstrate the visual calendar in the app: blue for slots that need volunteers, green for dates already selected, grey for holidays or locked dates.\n"
        "3. Emphasize that picking early gives them their preferred dates. Once the window closes, admins manually assign remaining slots.")

    # ---------------------------------------------------------
    # SLIDE 6: Task 3 - Recording Student Violations
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_light_slate)
    add_slide_header(slide, "Task 3: Logging Student Violations", "Campus Discipline")

    # Text instructions
    violation_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    v_tf = violation_text.text_frame
    v_tf.word_wrap = True
    v_tf.margin_left = v_tf.margin_top = v_tf.margin_right = v_tf.margin_bottom = 0
    
    p = v_tf.paragraphs[0]
    p.text = "Step-by-step violation recording:"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark_navy
    p.space_after = Pt(12)

    bullets_t3 = [
        "1. Active Duty Required: You must be checked in to record a student violation. The system verifies this state automatically.",
        "2. Start Form: Tap 'Record Student Violation' on your Dashboard.",
        "3. Select Student: Type the name or registration number. The database will auto-suggest matching records.",
        "4. Choose Violation: Pick from the list (e.g., 'Late to Assembly', 'Dress Code'). The preconfigured fine amount will load.",
        "5. Save Details: Set to 'Warning Only' if appropriate, add relevant notes, and tap 'Save Violation'."
    ]
    for b in bullets_t3:
        bp = v_tf.add_paragraph()
        bp.text = b
        bp.font.name = "Arial"
        bp.font.size = Pt(13)
        bp.font.color.rgb = c_text_ink
        bp.space_after = Pt(10)

    # Quick Add Mode box
    qa_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.1), Inches(6.5), Inches(1.1))
    qa_box.fill.solid()
    qa_box.fill.fore_color.rgb = RGBColor(236, 253, 245) # light green
    qa_box.line.color.rgb = c_accent_green
    qa_box.line.width = Pt(1)
    qtf = qa_box.text_frame
    qtf.word_wrap = True
    qtf.margin_top = qtf.margin_bottom = qtf.margin_left = qtf.margin_right = Inches(0.12)
    qp = qtf.paragraphs[0]
    qp.text = "💡 Quick-add mode:\nLogging a group? Enable 'Quick-Add Mode'. The screen will stay open and preserve your selected violation type, letting you search and log multiple students in quick succession."
    qp.font.name = "Arial"
    qp.font.size = Pt(10.5)
    qp.font.color.rgb = RGBColor(6, 95, 70) # dark green

    # Card graphic right side
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = c_white
    card.line.color.rgb = c_border_gray
    card.line.width = Pt(1)
    
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.margin_top = ctf.margin_left = ctf.margin_right = Inches(0.3)
    
    cp = ctf.paragraphs[0]
    cp.text = "RECORD STUDENT VIOLATION"
    cp.font.name = "Arial"
    cp.font.size = Pt(12)
    cp.font.bold = True
    cp.font.color.rgb = c_text_muted
    cp.space_after = Pt(14)

    fields = [
        "STUDENT SEARCH",
        "Aarav Menon — 21PH045",
        "VIOLATION TYPE",
        "Late to assembly area (₹50 fine)"
    ]
    for idx, f in enumerate(fields):
        fp = ctf.add_paragraph()
        fp.text = f
        if idx % 2 == 0:
            fp.font.name = "Courier New"
            fp.font.size = Pt(9.5)
            fp.font.bold = True
            fp.font.color.rgb = c_text_muted
            fp.space_after = Pt(2)
        else:
            fp.font.name = "Arial"
            fp.font.size = Pt(13)
            fp.font.bold = True
            fp.font.color.rgb = c_text_ink
            fp.space_after = Pt(12)

    # Submit button simulation
    s_btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(4.8), Inches(3.9), Inches(0.5))
    s_btn.fill.solid()
    s_btn.fill.fore_color.rgb = c_brand_blue
    s_btn.line.fill.background()
    s_tf = s_btn.text_frame
    sp = s_tf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.text = "Log Violation"
    sp.font.name = "Arial"
    sp.font.size = Pt(12)
    sp.font.bold = True
    sp.font.color.rgb = c_white

    add_notes(slide,
        "Speaker Notes: Describe how to register disciplinary violations.\n\n"
        "Key points:\n"
        "1. Clarify that the system requires check-in to avoid ad-hoc reports outside duty hours.\n"
        "2. Explain that the list of students is dynamically fetched from the database, updated regularly via administrative Excel uploads.\n"
        "3. Highlight that fines are strictly preset by systemic policy, removing arguments or discrepancies about how much to fine a student.")

    # ---------------------------------------------------------
    # SLIDE 7: Task 4 - Requesting Swaps
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_light_slate)
    add_slide_header(slide, "Task 4: Swapping Duties / Reassignment", "Schedule Management")

    # Text instructions
    swap_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    sw_tf = swap_text.text_frame
    sw_tf.word_wrap = True
    sw_tf.margin_left = sw_tf.margin_top = sw_tf.margin_right = sw_tf.margin_bottom = 0
    
    p = sw_tf.paragraphs[0]
    p.text = "When you cannot attend an assigned duty:"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark_navy
    p.space_after = Pt(12)

    bullets_t4 = [
        "1. Open the Dashboard and find the list under 'Upcoming Duties'.",
        "2. Tap the '🔄 Request Reassignment' button next to the relevant slot.",
        "3. Select a Colleague: Pick an eligible faculty member from the dropdown list. (Only active faculty members are shown).",
        "4. Provide a Reason: Input a brief explanation (e.g., 'Health issue', 'Exam duty') so your colleague understands.",
        "5. Tap 'Send Request'. The duty stays pending in your upcoming list."
    ]
    for b in bullets_t4:
        bp = sw_tf.add_paragraph()
        bp.text = b
        bp.font.name = "Arial"
        bp.font.size = Pt(13.5)
        bp.font.color.rgb = c_text_ink
        bp.space_after = Pt(11)

    # Info callout
    callout = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.1), Inches(6.5), Inches(1.2))
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(254, 243, 199) # light amber
    callout.line.color.rgb = c_accent_amber
    callout.line.width = Pt(1)
    ctf = callout.text_frame
    ctf.word_wrap = True
    ctf.margin_top = ctf.margin_bottom = ctf.margin_left = ctf.margin_right = Inches(0.12)
    cp = ctf.paragraphs[0]
    cp.text = "⚠️ Crucial rule:\nThe duty is NOT transferred until your colleague explicitly accepts your request. Until they click 'Accept', you are still officially responsible for that duty slot."
    cp.font.name = "Arial"
    cp.font.size = Pt(11)
    cp.font.bold = True
    cp.font.color.rgb = RGBColor(146, 64, 14)

    # Colleague card mockup
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = c_white
    card.line.color.rgb = c_border_gray
    card.line.width = Pt(1)
    
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.margin_top = ctf.margin_left = ctf.margin_right = Inches(0.3)
    
    cp = ctf.paragraphs[0]
    cp.text = "REQUEST REASSIGNMENT"
    cp.font.name = "Arial"
    cp.font.size = Pt(12)
    cp.font.bold = True
    cp.font.color.rgb = c_text_muted
    cp.space_after = Pt(14)

    fields = [
        "DUTY SLOT TO TRANSFER",
        "Saturday, 15 March · Afternoon",
        "SELECT COLLEAGUE",
        "Dr. R. Nair (Department of Pharmacology)",
        "REASON FOR REQUEST",
        "Attending university seminar."
    ]
    for idx, f in enumerate(fields):
        fp = ctf.add_paragraph()
        fp.text = f
        if idx % 2 == 0:
            fp.font.name = "Courier New"
            fp.font.size = Pt(9.5)
            fp.font.bold = True
            fp.font.color.rgb = c_text_muted
            fp.space_after = Pt(2)
        else:
            fp.font.name = "Arial"
            fp.font.size = Pt(12.5)
            fp.font.bold = True
            fp.font.color.rgb = c_text_ink
            fp.space_after = Pt(10)

    # Submit button simulation
    s_btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(4.9), Inches(3.9), Inches(0.45))
    s_btn.fill.solid()
    s_btn.fill.fore_color.rgb = c_brand_blue
    s_btn.line.fill.background()
    s_tf = s_btn.text_frame
    sp = s_tf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.text = "Send Swap Request"
    sp.font.name = "Arial"
    sp.font.size = Pt(12)
    sp.font.bold = True
    sp.font.color.rgb = c_white

    add_notes(slide,
        "Speaker Notes: Walk through duty reassignment.\n\n"
        "Key points:\n"
        "1. Faculty swaps are regular occurrences. The app handles this digitally to avoid manual coordinator overrides.\n"
        "2. Highlight that the colleague must be chosen from the active eligibility list (the app handles schedule exclusions automatically).\n"
        "3. Emphasize the accountability rule: they cannot simply abandon a slot; they must ensure the swap is confirmed and accepted.")

    # ---------------------------------------------------------
    # SLIDE 8: Task 5 - Responding to Swaps
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_light_slate)
    add_slide_header(slide, "Task 5: Answering Cover Swap Requests", "Incoming Tasks")

    # Text instructions
    in_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    itf = in_text.text_frame
    itf.word_wrap = True
    itf.margin_left = itf.margin_top = itf.margin_right = itf.margin_bottom = 0
    
    p = itf.paragraphs[0]
    p.text = "When a colleague asks you to cover:"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark_navy
    p.space_after = Pt(12)

    bullets_t5 = [
        "1. You will receive an instant notification directly in your Telegram chat.",
        "2. Open the app. Incoming requests will be listed near the top of the Dashboard under 'Reassignment Requests'.",
        "3. Check Details: Review the date, session type, requesting colleague name, and their provided reason.",
        "4. Accept: Tap the green 'Accept' button to take over the duty. The slot instantly moves to your schedule.",
        "5. Decline: Tap 'Reject' if you cannot make it. The duty remains with the original requester."
    ]
    for b in bullets_t5:
        bp = itf.add_paragraph()
        bp.text = b
        bp.font.name = "Arial"
        bp.font.size = Pt(13.5)
        bp.font.color.rgb = c_text_ink
        bp.space_after = Pt(12)

    # Info callout
    info_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.1), Inches(6.5), Inches(1.1))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = RGBColor(239, 246, 255) # light blue
    info_box.line.color.rgb = c_brand_blue
    info_box.line.width = Pt(1)
    i_tf = info_box.text_frame
    i_tf.word_wrap = True
    i_tf.margin_top = i_tf.margin_bottom = i_tf.margin_left = i_tf.margin_right = Inches(0.12)
    ip = i_tf.paragraphs[0]
    ip.text = "⚡ Telegram Shortcuts:\nYou do not even need to open the app! The Telegram bot notification contains interactive 'Accept' and 'Decline' buttons. Tapping those handles the swap instantly."
    ip.font.name = "Arial"
    ip.font.size = Pt(11)
    ip.font.color.rgb = c_brand_blue

    # Card graphic right side (Telegram Bot message simulation)
    bot_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.5))
    bot_card.fill.solid()
    bot_card.fill.fore_color.rgb = RGBColor(30, 41, 59) # dark theme mimic for bot
    bot_card.line.color.rgb = c_brand_blue
    bot_card.line.width = Pt(2)
    
    b_tf = bot_card.text_frame
    b_tf.word_wrap = True
    b_tf.margin_top = b_tf.margin_left = b_tf.margin_right = Inches(0.3)
    
    bp1 = b_tf.paragraphs[0]
    bp1.text = "💬 TELEGRAM BOT NOTIFICATION"
    bp1.font.name = "Courier New"
    bp1.font.size = Pt(10)
    bp1.font.bold = True
    bp1.font.color.rgb = RGBColor(147, 197, 253) # light blue
    bp1.space_after = Pt(14)

    bp2 = b_tf.add_paragraph()
    bp2.text = "🔄 Cover Request from Priya S.:\n\n- Date: Saturday, 15 March\n- Session: Afternoon (1:00 PM)\n- Reason: Personal commitment.\n\nWould you like to accept this shift?"
    bp2.font.name = "Arial"
    bp2.font.size = Pt(12)
    bp2.font.color.rgb = c_white
    bp2.space_after = Pt(20)

    # Green button
    btn_g = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(4.8), Inches(1.9), Inches(0.45))
    btn_g.fill.solid()
    btn_g.fill.fore_color.rgb = c_accent_green
    btn_g.line.fill.background()
    g_tf = btn_g.text_frame
    gp = g_tf.paragraphs[0]
    gp.alignment = PP_ALIGN.CENTER
    gp.text = "Accept"
    gp.font.name = "Arial"
    gp.font.size = Pt(11.5)
    gp.font.bold = True
    gp.font.color.rgb = c_white

    # Red button
    btn_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(4.8), Inches(1.9), Inches(0.45))
    btn_r.fill.solid()
    btn_r.fill.fore_color.rgb = c_accent_red
    btn_r.line.fill.background()
    r_tf = btn_r.text_frame
    rp = r_tf.paragraphs[0]
    rp.alignment = PP_ALIGN.CENTER
    rp.text = "Decline"
    rp.font.name = "Arial"
    rp.font.size = Pt(11.5)
    rp.font.bold = True
    rp.font.color.rgb = c_white

    add_notes(slide,
        "Speaker Notes: Describe what happens when a swap request is received.\n\n"
        "Key points:\n"
        "1. Highlight how convenient it is to click 'Accept' or 'Decline' directly inside Telegram chat without needing to log in to the web browser.\n"
        "2. Advise them to respond promptly so colleagues have time to search for other options if declined.\n"
        "3. Mention that once accepted, the scheduling system sends a confirmation notification to both users.")

    # ---------------------------------------------------------
    # SLIDE 9: Tracking & Self-Management
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_light_slate)
    add_slide_header(slide, "Track Your Record & Fix Mistakes", "Self-Management")

    # Text instructions
    track_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    t_tf = track_text.text_frame
    t_tf.word_wrap = True
    t_tf.margin_left = t_tf.margin_top = t_tf.margin_right = t_tf.margin_bottom = 0
    
    p = t_tf.paragraphs[0]
    p.text = "Managing your history:"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark_navy
    p.space_after = Pt(12)

    bullets_t6 = [
        "1. Attendance Summary: Open the 'Attendance' tab to view a monthly summary of your check-in statistics, including late arrivals and missed checkouts.",
        "2. Logged Violations: Open the 'Violations' tab to see every student record you have created, categorized by date and status.",
        "3. Correcting Mistakes: If you log a student by mistake, you can tap the 'Delete' button directly on that violation card.",
        "4. Flagging for Admin: If you are unsure of a record (e.g., student disputes a fine), tap 'Flag for Review'. This marks the violation and submits a alert to the Admin dashboard."
    ]
    for b in bullets_t6:
        bp = t_tf.add_paragraph()
        bp.text = b
        bp.font.name = "Arial"
        bp.font.size = Pt(13)
        bp.font.color.rgb = c_text_ink
        bp.space_after = Pt(10)

    # Warning callout
    del_warn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.1), Inches(6.5), Inches(1.1))
    del_warn.fill.solid()
    del_warn.fill.fore_color.rgb = RGBColor(254, 242, 242) # light red
    del_warn.line.color.rgb = c_accent_red
    del_warn.line.width = Pt(1)
    dtf = del_warn.text_frame
    dtf.word_wrap = True
    dtf.margin_top = dtf.margin_bottom = dtf.margin_left = dtf.margin_right = Inches(0.12)
    dp = dtf.paragraphs[0]
    dp.text = "⚠️ Delete constraints:\nViolations can only be deleted by the recording faculty while they are 'active'. Once an admin reviews and locks the report, edits or deletes require administrator override."
    dp.font.name = "Arial"
    dp.font.size = Pt(10.5)
    dp.font.color.rgb = RGBColor(153, 27, 27)

    # Graphic: Mock Stats Box (Right Side)
    stats_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.5))
    stats_card.fill.solid()
    stats_card.fill.fore_color.rgb = c_white
    stats_card.line.color.rgb = c_border_gray
    stats_card.line.width = Pt(1)
    
    stf = stats_card.text_frame
    stf.word_wrap = True
    stf.margin_top = stf.margin_left = stf.margin_right = Inches(0.3)
    
    sp = stf.paragraphs[0]
    sp.text = "MY ATTENDANCE SUMMARY"
    sp.font.name = "Arial"
    sp.font.size = Pt(12)
    sp.font.bold = True
    sp.font.color.rgb = c_text_muted
    sp.space_after = Pt(20)

    stats = [
        "On-Time Check Ins", "8 sessions", c_accent_green,
        "Late Arrivals", "1 session", c_accent_amber,
        "Forgot Check Out (Auto)", "2 sessions", c_accent_red
    ]
    for idx in range(0, len(stats), 3):
        label = stats[idx]
        val = stats[idx+1]
        color = stats[idx+2]
        
        stat_p = stf.add_paragraph()
        stat_p.text = f"• {label}: "
        stat_p.font.name = "Arial"
        stat_p.font.size = Pt(13.5)
        stat_p.font.color.rgb = c_text_ink
        
        # Append styled count
        run = stat_p.add_run()
        run.text = val
        run.font.bold = True
        run.font.color.rgb = color
        
        stat_p.space_after = Pt(16)

    add_notes(slide,
        "Speaker Notes: Introduce self-management and audits.\n\n"
        "Key points:\n"
        "1. The Attendance and Violations tabs give faculty full transparency over what they have logged.\n"
        "2. Emphasize that mistakes happen. The system allows faculty to delete or edit records directly before audits, preserving database clean states.\n"
        "3. Explain the difference: Delete removes it immediately, while Flag keeps the record in the view but highlights it in red on the Admin panel for discussion.")

    # ---------------------------------------------------------
    # SLIDE 10: Telegram Reminders & System Help
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_light_slate)
    add_slide_header(slide, "Telegram Reminders & Getting Help", "Support & Communications")

    # Text instructions
    help_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    htf = help_text.text_frame
    htf.word_wrap = True
    htf.margin_left = htf.margin_top = htf.margin_right = htf.margin_bottom = 0
    
    p = htf.paragraphs[0]
    p.text = "How Telegram Bot assists you:"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark_navy
    p.space_after = Pt(10)

    bullets_t7 = [
        "• Duty Reminder: Sent 30 minutes before your scheduled morning or afternoon session starts.",
        "• Clock-out Reminder: Placed 15 minutes before the session cutoff hour so you check out manually.",
        "• Swap Alerts: Delivers cover swap notifications with instant action buttons.",
        "• Command Shortcuts: Message '/menu' to check your active slots, next duty date, and window status dynamically."
    ]
    for b in bullets_t7:
        bp = htf.add_paragraph()
        bp.text = b
        bp.font.name = "Arial"
        bp.font.size = Pt(13)
        bp.font.color.rgb = c_text_ink
        bp.space_after = Pt(10)

    # Help block
    contact_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.7), Inches(6.5), Inches(1.6))
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = RGBColor(254, 243, 199) # light amber
    contact_box.line.color.rgb = c_accent_amber
    contact_box.line.width = Pt(1)
    ctf2 = contact_box.text_frame
    ctf2.word_wrap = True
    ctf2.margin_top = ctf2.margin_bottom = ctf2.margin_left = ctf2.margin_right = Inches(0.15)
    cp2 = ctf2.paragraphs[0]
    cp2.text = "🆘 NEED TECHNICAL SUPPORT?\nIf you experience any of the following, contact your Admin immediately:\n1. Your OTP code does not arrive in Telegram.\n2. You are locked out due to multiple failed OTP login attempts.\n3. You need to drop a picked slot after the schedule window has locked."
    cp2.font.name = "Arial"
    cp2.font.size = Pt(11)
    cp2.font.bold = True
    cp2.font.color.rgb = RGBColor(146, 64, 14)

    # Visual representation: Bot screen mock (Right column)
    bot_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.5))
    bot_card.fill.solid()
    bot_card.fill.fore_color.rgb = c_white
    bot_card.line.color.rgb = c_border_gray
    bot_card.line.width = Pt(1)
    
    b_tf2 = bot_card.text_frame
    b_tf2.word_wrap = True
    b_tf2.margin_top = b_tf2.margin_left = b_tf2.margin_right = Inches(0.3)
    
    bp1 = b_tf2.paragraphs[0]
    bp1.text = "📋 TELEGRAM SHORTCUTS (/menu)"
    bp1.font.name = "Arial"
    bp1.font.size = Pt(12)
    bp1.font.bold = True
    bp1.font.color.rgb = c_brand_blue
    bp1.space_after = Pt(14)

    # Button list mock
    b_texts = [
        "📋 My Duty Slots",
        "⏭️ Next Duty",
        "📅 Scheduling Window Status"
    ]
    b_top = Inches(2.7)
    for txt in b_texts:
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), b_top, Inches(3.9), Inches(0.48))
        btn.fill.solid()
        btn.fill.fore_color.rgb = RGBColor(241, 245, 249)
        btn.line.color.rgb = c_border_gray
        btn.line.width = Pt(1)
        
        btf = btn.text_frame
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        bp.text = txt
        bp.font.name = "Arial"
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = c_text_ink
        b_top += Inches(0.65)

    add_notes(slide,
        "Speaker Notes: Outline bot actions and system support.\n\n"
        "Key points:\n"
        "1. Reinforce that Telegram is the primary alert mechanism. Advise them not to mute the bot during semesters.\n"
        "2. Introduce the '/menu' feature for quick checks on dates without opening a browser.\n"
        "3. Provide support guidelines. Clarify that while the app runs autonomously, admins hold full capability to reset user logins and manual overrides.")

    # ---------------------------------------------------------
    # SLIDE 11: Cheat Sheet Summary (Dark Theme)
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, c_dark_navy)
    
    # Accent decorative block
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c_brand_blue
    shape.line.fill.background()

    # Title Box
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.8), Inches(11.0), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "SIMS DMS — Faculty Quick Cheat Sheet"
    p.font.name = "Arial"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.space_after = Pt(4)
    
    sp = tf.add_paragraph()
    sp.text = "Keep this overview on hand during the semester."
    sp.font.name = "Arial"
    sp.font.size = Pt(14)
    sp.font.color.rgb = c_text_muted

    # 3-column summary grid
    col_data = [
        ("🔑 SIGN IN", [
            "1. Enter email in app.",
            "2. Get OTP on Telegram.",
            "3. Authenticate instantly."
        ], Inches(1.2)),
        ("📋 DUTY ROUTINE", [
            "1. Check In when arriving.",
            "2. Check Out when leaving.",
            "3. Clock out before auto-checkout."
        ], Inches(5.0)),
        ("🔄 SWAPPING SLOTS", [
            "1. Tap Reassignment request.",
            "2. Select eligible colleague.",
            "3. Pending until they accept."
        ], Inches(8.8))
    ]

    for title, points, left in col_data:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), Inches(3.4), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.color.rgb = c_brand_blue
        box.line.width = Pt(1.5)
        
        btf = box.text_frame
        btf.word_wrap = True
        btf.margin_top = btf.margin_left = btf.margin_right = Inches(0.25)
        
        tp = btf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tp.text = title
        tp.font.name = "Arial"
        tp.font.size = Pt(13)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(147, 197, 253)
        tp.space_after = Pt(16)
        
        for pt in points:
            pp = btf.add_paragraph()
            pp.text = pt
            pp.font.name = "Arial"
            pp.font.size = Pt(11.5)
            pp.font.color.rgb = c_white
            pp.space_after = Pt(14)

    add_notes(slide,
        "Speaker Notes: Concluding slide.\n\n"
        "Key points:\n"
        "1. Summarize the three pillars: logging in, logging attendance, and request reassignments.\n"
        "2. Hand out the printed one-page cheat sheet compiled from this slide deck.\n"
        "3. Thank the faculty for their cooperation and open the floor to questions.")

    # Save presentation
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation(r"c:\Users\sikha\Music\sims disclipne\SIMS-DMS-Faculty-Presentation.pptx")
